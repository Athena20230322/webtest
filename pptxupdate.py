from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE  # 添加这行导入


def set_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_logo(slide, img_path, position):
    """添加公司logo"""
    left, top = position
    slide.shapes.add_picture(img_path, left, top, height=Cm(1.5))


def create_modern_title_slide(prs, title_text, subtitle_text):
    """创建现代风格封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局

    # 设置渐变背景
    set_background(slide, (13, 43, 69))  # 深蓝色

    # 添加标题
    title = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(20), Cm(3))
    tf = title.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    p.font.bold = True

    # 添加副标题
    subtitle = slide.shapes.add_textbox(Cm(2), Cm(6), Cm(20), Cm(3))
    tf = subtitle.text_frame
    p = tf.add_paragraph()
    p.text = subtitle_text
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)  # 浅灰色

    # 添加装饰线（修正了这里的MSO_SHAPE使用）
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), Cm(5.5), Cm(5), Cm(0.2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 176, 240)  # 蓝色

    return slide


def create_content_slide(prs, title_text, content_items, layout='two_column'):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容布局
    set_background(slide, (240, 240, 240))  # 浅灰背景

    # 设置标题
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(13, 43, 69)

    # 添加内容
    content = slide.placeholders[1]
    tf = content.text_frame

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)

    return slide


def main():
    prs = Presentation()

    # 1. 封面页
    create_modern_title_slide(
        prs,
        "iCP APP及API測試上線報告",
        "報告人: Adan\n日期: 2025年7月\nicash Pay 上線前測試階段"
    )

    # 2. APP测试部分
    app_auto_content = [
        "• 自动化测试工具: Python + pytest + BlueStack模拟器",
        "• 测试范围: 功能回归测试、性能基准测试",
        "• 产出报告: 自动生成HTML报告"
    ]

    app_manual_content = [
        "• 手动测试框架: Node.js + Express.js",
        "• 测试重点:",
        "  - 反扫扣款流程",
        "  - 现金储值验证",
        "  - 授权绑定测试",
        "  - URL跳转检查",
        "  - WebQRCode扫描"
    ]

    create_content_slide(prs, "APP自动化测试", app_auto_content)
    create_content_slide(prs, "APP手动测试", app_manual_content)

    # 3. API测试部分
    api_auto_content = [
        "• 工具迁移: 从C#改为Python执行",
        "• 核心脚本: python run_icplogin.py",
        "• 测试类型: 接口契约测试、压力测试"
    ]

    api_manual_content = [
        "• Mock服务: 从C#改为Python实现",
        "• 测试工具: python api_tester.py",
        "• 辅助工具: Node.js测试小工具开发"
    ]

    create_content_slide(prs, "API自动化测试", api_auto_content)
    create_content_slide(prs, "API手动测试", api_manual_content)

    # 4. 测试报告页
    report_content = [
        "• APP自动化测试报告",
        "• APP手动测试报告",
        "• API自动化测试报告",
        "• API手动测试报告(一)",
        "• API手动测试报告(二)"
    ]

    create_content_slide(prs, "测试产出报告", report_content)

    # 保存文件
    prs.save('icp_test_report_美化版.pptx')
    print("PPTX美化完成: icp_test_report_美化版.pptx")


if __name__ == "__main__":
    main()